import xlwings as xw
from pathlib import Path
from pptx import Presentation
import comtypes.client



def main():
    # Path settings
    current_dir = Path(__file__).parent
    pptx_template_path = current_dir / "adt.pptx"

    # Conexión a Excel
    wb = xw.Book.caller()
    sht_panel = wb.sheets["PANEL"]
    context = sht_panel.range("B2").options(dict, expand="table", numbers=int).value
    folder_name = sht_panel.range("C13").value
    mes_name = sht_panel.range("C4").value
    # Iniciar plantilla
    pptx_template = Presentation(pptx_template_path)

    # Reemplazar valores de Excel en la plantilla de PowerPoint 
    PROYECTO = context["PROYECTO"]
    NOMBRE_CLIENTE = context["NOMBRE_CLIENTE"]
    MES = context["MES"]
    LUGAR = context["LUGAR"]
    FECHA = context["FECHA"]
    HORA = context["HORA"]
    HORA2 = context["HORA2"]
    A1 = context["A1"]
    A1 = str(A1)
    ASI1 = context["ASI1"]
    ASI1 = str(ASI1)
    A2 = context["A2"]
    A2 = str(A2)
    ASI2 = context["ASI2"]
    ASI2 = str(ASI2)
    for slide in pptx_template.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shape.text = shape.text.replace("{{ PROYECTO }}", PROYECTO)
                        shape.text = shape.text.replace("{{ NOMBRE_CLIENTE }}", NOMBRE_CLIENTE)
                        shape.text = shape.text.replace("{{ MES }}", MES)
                        shape.text = shape.text.replace("{{ LUGAR }}", LUGAR)
                        shape.text = shape.text.replace("{{ FECHA }}", FECHA)
                        shape.text = shape.text.replace("{{ HORA }}", HORA)
                        shape.text = shape.text.replace("{{ HORA2 }}", HORA2)
                        shape.text = shape.text.replace("{{ A1 }}", A1)
                        shape.text = shape.text.replace("{{ A2 }}", A2)
                        shape.text = shape.text.replace("{{ ASI1 }}", ASI1)
                        shape.text = shape.text.replace("{{ ASI2 }}", ASI2)


    output_pptx_name = f"C:\\Users\\polit\\OneDrive - HRM\\PROYECTOS\\2023\\{folder_name}\\PROYECTOS\\{mes_name}\\{context['PROYECTO']}.pptx"
    pptx_template.save(output_pptx_name)

    output_pdf = f"C:\\Users\\polit\\OneDrive - HRM\\PROYECTOS\\2023\\{folder_name}\\PROYECTOS\\{mes_name}\\{context['PROYECTO']}.pdf"
    input_ppt = str(output_pptx_name)
    output_pdf = str(output_pdf)
   
    
    ppt_to_pdf(input_ppt, output_pdf, wb)
def ppt_to_pdf(input_ppt, output_pdf, wb):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    ppt = powerpoint.Presentations.Open(input_ppt)
    ppt.SaveAs(output_pdf, FileFormat=32)
    ppt.Close()

    powerpoint.Quit()  


    # Show Message Box
    show_msgbox = wb.macro("Module1.ShowMsgBox")
    show_msgbox("Listo!")


if __name__ == "__main__":
        xw.Book("pptx_automation.xlsm").set_mock_caller()
        main()
